"""Slide preview images for the PPT library.

Windows uses WPS/PowerPoint to export real slide screenshots. Linux uses the
HTML renderer. A lightweight python-pptx fallback is last resort only.
"""

from __future__ import annotations

import io
import logging
import os
import threading
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont

from utils.cjk_fonts import ensure_cjk_preview_font

LOGGER = logging.getLogger(__name__)

PREVIEW_WIDTH = 1280
PREVIEW_HEIGHT = 720
_OFFICE_LOCK = threading.Lock()


def count_pptx_slides(pptx_path: str) -> int:
    from pptx import Presentation

    return len(Presentation(pptx_path).slides)


def render_pptx_via_office(pptx_path: str, dest_dir: str) -> list[str]:
    """Export slides with WPS or PowerPoint. Empty on Linux or if Office is missing."""
    if os.name != "nt":
        return []
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        LOGGER.info("[ppt.library] win32com is not installed; skip Office export")
        return []

    os.makedirs(dest_dir, exist_ok=True)
    abs_pptx = os.path.abspath(pptx_path)
    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        with _OFFICE_LOCK:
            for prog_id in ("Kwpp.Application", "PowerPoint.Application", "WPP.Application"):
                try:
                    app = win32com.client.Dispatch(prog_id)
                    LOGGER.info("[ppt.library] Office export via %s", prog_id)
                    break
                except Exception:
                    app = None
            if app is None:
                return []
            for attr, value in (("Visible", False), ("DisplayAlerts", 0)):
                try:
                    setattr(app, attr, value)
                except Exception:
                    pass
            try:
                presentation = app.Presentations.Open(abs_pptx, True, False, False)
            except Exception:
                presentation = app.Presentations.Open(abs_pptx, WithWindow=False)
            paths: list[str] = []
            count = int(presentation.Slides.Count)
            for index in range(1, count + 1):
                dest = os.path.abspath(os.path.join(dest_dir, f"slide_{index}.png"))
                presentation.Slides(index).Export(dest, "PNG", PREVIEW_WIDTH, PREVIEW_HEIGHT)
                if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                    paths.append(dest)
            return paths
    except Exception:
        LOGGER.exception("[ppt.library] Office slide export failed path=%s", pptx_path)
        return []
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def render_pptx_slide_previews(pptx_path: str, dest_dir: str) -> list[str]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.makedirs(dest_dir, exist_ok=True)
    presentation = Presentation(pptx_path)
    slide_w = int(presentation.slide_width or 1)
    slide_h = int(presentation.slide_height or 1)
    font = _load_cjk_font(36)
    small_font = _load_cjk_font(22)
    paths: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        canvas = Image.new("RGB", (PREVIEW_WIDTH, PREVIEW_HEIGHT), _slide_bg_rgb(slide))
        draw = ImageDraw.Draw(canvas)
        for shape in _iter_shapes(slide.shapes):
            box = _shape_box(shape, slide_w, slide_h)
            if not box:
                continue
            shape_type = getattr(shape, "shape_type", None)
            if shape_type in {MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM}:
                color = _shape_fill_rgb(shape)
                if color:
                    radius = min(28, max(0, (box[2] - box[0]) // 10))
                    draw.rounded_rectangle(box, radius=radius, fill=color)
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                _paste_picture(canvas, shape, slide_w, slide_h)
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    use_font = font if (box[3] - box[1]) > 48 else small_font
                    draw.text(
                        (box[0] + 8, box[1] + 6),
                        text.replace("\x0b", "\n")[:80],
                        font=use_font,
                        fill=(31, 22, 59),
                    )
        dest = os.path.join(dest_dir, f"slide_{index}.png")
        canvas.save(dest, "PNG", optimize=True)
        paths.append(dest)
    return paths


def _iter_shapes(shapes) -> Iterable:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            nested = getattr(shape, "shapes", None)
            if nested is not None:
                yield from _iter_shapes(nested)
            continue
        yield shape


def _slide_bg_rgb(slide) -> tuple[int, int, int]:
    try:
        from pptx.enum.dml import MSO_FILL

        fill = slide.background.fill
        if fill.type == MSO_FILL.SOLID:
            rgb = fill.fore_color.rgb
            return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except Exception:
        pass
    return (247, 248, 251)


def _shape_fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        from pptx.enum.dml import MSO_FILL

        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            rgb = fill.fore_color.rgb
            return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except Exception:
        return None
    return None


def _emu_to_px(value, slide_emu: int, preview: int) -> int:
    try:
        return int(int(value) * preview / max(slide_emu, 1))
    except Exception:
        return 0


def _shape_box(shape, slide_w: int, slide_h: int) -> tuple[int, int, int, int] | None:
    x = _emu_to_px(getattr(shape, "left", 0) or 0, slide_w, PREVIEW_WIDTH)
    y = _emu_to_px(getattr(shape, "top", 0) or 0, slide_h, PREVIEW_HEIGHT)
    width = max(1, _emu_to_px(getattr(shape, "width", 0) or 0, slide_w, PREVIEW_WIDTH))
    height = max(1, _emu_to_px(getattr(shape, "height", 0) or 0, slide_h, PREVIEW_HEIGHT))
    if width <= 1 and height <= 1:
        return None
    return (x, y, x + width, y + height)


def _paste_picture(canvas: Image.Image, shape, slide_w: int, slide_h: int) -> bool:
    try:
        blob = shape.image.blob
        picture = Image.open(io.BytesIO(blob)).convert("RGBA")
    except Exception:
        return False
    x = _emu_to_px(getattr(shape, "left", 0) or 0, slide_w, PREVIEW_WIDTH)
    y = _emu_to_px(getattr(shape, "top", 0) or 0, slide_h, PREVIEW_HEIGHT)
    width = max(1, _emu_to_px(getattr(shape, "width", 0) or 0, slide_w, PREVIEW_WIDTH))
    height = max(1, _emu_to_px(getattr(shape, "height", 0) or 0, slide_h, PREVIEW_HEIGHT))
    try:
        picture = picture.resize((width, height), Image.Resampling.LANCZOS)
        canvas.paste(picture, (x, y), picture)
        return True
    except Exception:
        LOGGER.debug("[ppt.library] skip picture on slide preview", exc_info=True)
        return False


def _load_cjk_font(size: int):
    font_path = ensure_cjk_preview_font()
    if font_path:
        try:
            return ImageFont.truetype(font_path, size=size)
        except OSError:
            pass
    return ImageFont.load_default()
